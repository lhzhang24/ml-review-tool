"""
ML 速查工坊 — 课件上传 → AI 解析 → 公式速查卡 + 重点整理
支持格式: PDF / PPTX / 图片 (PNG, JPG, WEBP)
"""

import os
import re
import json
import base64
import io
import hashlib
import zipfile
import tempfile
import traceback
from pathlib import Path
from datetime import datetime

from flask import Flask, request, jsonify, render_template, send_from_directory
from werkzeug.utils import secure_filename

# ── 文件提取 ────────────────────────────────────────────
try:
    from PyPDF2 import PdfReader
except ImportError:
    PdfReader = None

try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from PIL import Image
except ImportError:
    Image = None

try:
    import pytesseract
except ImportError:
    pytesseract = None

# ── App 初始化 ───────────────────────────────────────────
BASE_DIR = Path(__file__).resolve().parent
UPLOAD_DIR = BASE_DIR / "uploads"
OUTPUT_DIR = BASE_DIR / "outputs"
STATS_FILE = BASE_DIR / "stats.json"
UPLOAD_DIR.mkdir(exist_ok=True)
OUTPUT_DIR.mkdir(exist_ok=True)

# ── 统计持久化（雷达图数据）───────────────────────────────
def _load_stats():
    """读取 stats.json，不存在则返回默认值"""
    defaults = {
        "mastery": {"beginner": 0, "intermediate": 0, "god": 0},
        "urgency": {"rush": 0, "relaxed": 0, "giveup": 0},
    }
    if not STATS_FILE.exists():
        return defaults
    try:
        data = json.loads(STATS_FILE.read_text(encoding="utf-8"))
        # 补全缺失的键
        for key in defaults:
            if key not in data:
                data[key] = defaults[key]
            for subkey, default_val in defaults[key].items():
                if subkey not in data[key]:
                    data[key][subkey] = default_val
        return data
    except Exception:
        return defaults

def _save_stats(data: dict):
    """写入 stats.json"""
    try:
        STATS_FILE.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")
    except Exception:
        pass

def increment_stats(mastery: str, urgency: str):
    """递增对应选项的计数"""
    data = _load_stats()
    if mastery in data["mastery"]:
        data["mastery"][mastery] += 1
    if urgency in data["urgency"]:
        data["urgency"][urgency] += 1
    _save_stats(data)
    return data

app = Flask(__name__)
app.config["MAX_CONTENT_LENGTH"] = 500 * 1024 * 1024  # 500MB — 支持大压缩包

ALLOWED_EXTENSIONS = {"pdf", "pptx", "ppt", "png", "jpg", "jpeg", "webp", "bmp", "tiff", "zip"}

# ── 工具函数 ────────────────────────────────────────────

def allowed_file(filename):
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS

def safe_basename(original: str) -> str:
    """保留 Unicode（中文等）的文件名清理，只剔除危险字符"""
    import unicodedata
    # 取纯文件名（去掉路径前缀）
    name = os.path.basename(original)
    # 剔除控制字符和路径分隔符
    name = re.sub(r'[\x00-\x1f\x7f/\\:*?"<>|]', '_', name)
    # 去首尾空格和点
    name = name.strip(' .')
    if not name:
        # 空名 → 用哈希兜底
        name = f"file_{hashlib.md5(original.encode('utf-8', errors='replace')).hexdigest()[:8]}"
    return name

def get_ext(filename: str) -> str:
    """获取小写扩展名"""
    return filename.rsplit(".", 1)[1].lower() if "." in filename else ""

def extract_text_pdf(filepath: str) -> str:
    """提取 PDF 文本"""
    text = ""
    # 优先用 pdfplumber（效果更好）
    if pdfplumber:
        try:
            with pdfplumber.open(filepath) as pdf:
                for page in pdf.pages:
                    t = page.extract_text()
                    if t:
                        text += t + "\n"
            if text.strip():
                return text.strip()
        except Exception:
            pass
    # 回退到 PyPDF2
    if PdfReader:
        try:
            reader = PdfReader(filepath)
            for page in reader.pages:
                t = page.extract_text()
                if t:
                    text += t + "\n"
        except Exception:
            pass
    return text.strip()

def extract_text_pptx(filepath: str) -> str:
    """提取 PPTX 文本（也尝试处理 .ppt，但不保证成功）"""
    if not Presentation:
        return "[工具缺失] python-pptx 未安装"
    text = ""
    ext = filepath.rsplit(".", 1)[1].lower() if "." in filepath else ""
    try:
        prs = Presentation(filepath)
        for slide_num, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            ct = cell.text.strip()
                            if ct:
                                row_text.append(ct)
                        if row_text:
                            texts.append(" | ".join(row_text))
            if texts:
                text += f"── 第 {slide_num} 页 ──\n" + "\n".join(texts) + "\n\n"
    except Exception as e:
        if ext == "ppt":
            return ("[格式不支持] 旧版 .ppt 格式无法直接读取（python-pptx 只支持 .pptx）。\n"
                    "请用 PowerPoint 打开 → 另存为 → 选择 .pptx 格式，再重新上传。\n"
                    "—— 或者导出为 PDF 后上传 PDF。")
        return f"[解析失败] {str(e)}"
    if not text.strip():
        return "[提示] 未检测到文本，可能是纯图片 PPT，建议导出为 PDF 后上传。"
    return text.strip()

def extract_text_image(filepath: str) -> str:
    """从图片中提取文字（OCR）"""
    if not Image or not pytesseract:
        return "[提示] 图片 OCR 需要安装 tesseract: brew install tesseract tesseract-lang"
    try:
        img = Image.open(filepath)
        text = pytesseract.image_to_string(img, lang="chi_sim+eng")
        return text.strip()
    except Exception as e:
        return f"[OCR 失败] {str(e)}"

def extract_text(filepath: str) -> dict:
    """根据文件类型提取文本"""
    ext = filepath.rsplit(".", 1)[1].lower()
    if ext == "pdf":
        text = extract_text_pdf(filepath)
        ftype = "PDF"
    elif ext in ("pptx", "ppt"):
        text = extract_text_pptx(filepath)
        ftype = "PPT"
    else:
        text = extract_text_image(filepath)
        ftype = "图片"
    return {"text": text, "type": ftype, "filename": os.path.basename(filepath)}

def extract_zip(filepath: str) -> tuple:
    """解压 ZIP 并提取所有支持的课件文件。
    返回 (extracted_paths: list, warnings: list)"""
    extracted_paths = []
    warnings = []
    try:
        with zipfile.ZipFile(filepath, 'r') as zf:
            for name in zf.namelist():
                # 跳过目录
                if name.endswith('/'):
                    continue
                # 跳过 Mac 隐藏文件（__MACOSX/ 或 ._ 前缀）
                basename = os.path.basename(name)
                basename_lower = basename.lower()
                if '__macosx' in basename_lower or basename_lower.startswith('._'):
                    continue
                if basename.startswith('.'):
                    continue

                ext = get_ext(name)
                if ext not in ALLOWED_EXTENSIONS or ext == 'zip':
                    if ext:
                        warnings.append(f"跳过不支持的文件类型 .{ext}: {os.path.basename(name)}")
                    continue

                safe_name = safe_basename(name)
                dest = UPLOAD_DIR / safe_name
                # 避免重名覆盖
                counter = 1
                stem, real_ext = os.path.splitext(safe_name)
                while dest.exists():
                    dest = UPLOAD_DIR / f"{stem}_{counter}{real_ext}"
                    counter += 1
                with zf.open(name) as src, open(dest, 'wb') as dst:
                    dst.write(src.read())
                extracted_paths.append(str(dest))
    except zipfile.BadZipFile:
        warnings.append("ZIP 文件损坏，无法解压")
    except Exception as e:
        warnings.append(f"ZIP 解压异常: {str(e)[:200]}")
    return extracted_paths, warnings

def build_prompt(extracted_text: str, file_type: str, filename: str,
                 mastery: str = "intermediate", urgency: str = None) -> str:
    """构建 AI 提示词（根据掌握程度和紧迫程度调整风格）"""
    latex_example = r"$J(\theta)=\frac{1}{2m}\sum(h_\theta(x)-y)^2$"
    json_example_f = r'[{"formula": "LaTeX", "name": "名称", "meaning": "一句话含义", "source": "文件名"}]'
    json_example_p = r'[{"topic": "知识点", "summary": "核心思想", "priority": "★★★/★★/★", "pitfall": "易错提醒", "source": "文件名"}]'

    # 根据掌握程度调整输出风格
    mastery_instruction = {
        "beginner": (
            "【输出风格：零基础模式】\n"
            "- 每个公式都要用大白话解释含义，不要默认用户懂任何术语\n"
            "- 考点解释要详细，给出具体例子\n"
            "- 给出「这个公式是干嘛的」的一句话说明\n"
        ),
        "intermediate": (
            "【输出风格：进阶模式】\n"
            "- 公式解释简洁但完整，侧重使用场景和易错点\n"
            "- 考点突出考试重点和常见陷阱\n"
        ),
        "god": (
            "【输出风格：大佬模式】\n"
            "- 公式只给 LaTeX 和极简含义，不废话\n"
            "- 考点只列高频精华，skip 基础概念\n"
            "- 可以适量使用英文术语\n"
        ),
    }.get(mastery, "【输出风格：进阶模式】\n- 适中详细程度\n")

    # 根据时间紧迫程度调整策略
    urgency_instruction = {
        "rush": (
            "【复习策略：急速模式】\n"
            "- 只提取最高优先级的公式（与 ★★★ 考点相关的）\n"
            "- 每个公式/考点用最精炼的语言，直击要害\n"
            "- 可以跳过详细解释，给「秒背版」\n"
        ),
        "relaxed": (
            "【复习策略：从容模式】\n"
            "- 详细完整，每个公式给使用场景和推导思路\n"
            "- 考点配合例题思路或直观解释\n"
            "- 可以适当拓展相关知识点\n"
        ),
        "giveup": (
            "【复习策略：佛系陪伴模式】\n"
            "- 语气轻松幽默，适当使用网络梗和 emoji\n"
            "- 公式解释用最接地气的方式，帮你建立直觉\n"
            "- 考点标注「背这个就够了的保底内容」\n"
        ),
    }.get(urgency, "")

    parts = []
    parts.append(f"你是一位计算机专业的机器学习助教。请根据以下 {file_type} 课件内容，输出两样东西：\n")
    parts.append(mastery_instruction)
    if urgency_instruction:
        parts.append(urgency_instruction)
    parts.append("【一、公式速查卡】（A4 一页可打印）")
    parts.append(f"- 提取课件中出现的所有数学公式（用 LaTeX 格式，如 {latex_example}）")
    parts.append("- 每个公式标注含义（解决什么问题）")
    parts.append(f"- 每个公式的 source 字段填写来源文件名：{filename}")
    parts.append("- 按重要程度排序，最多 30 条，尽量多提取\n")
    parts.append("【二、重点考点整理】")
    parts.append("- 列出课件覆盖的知识点（每条一句话概括核心思想）")
    parts.append("- 标注哪些大概率是考试重点（★★★/★★/★）")
    parts.append("- 给出每个知识点的易错/易混淆提醒")
    parts.append(f"- 每个考点的 source 字段填写来源文件名：{filename}\n")
    parts.append("输出格式：")
    parts.append("```")
    parts.append(f"===FORMULA_CARD===\n{json_example_f}\n")
    parts.append(f"===KEY_POINTS===\n{json_example_p}")
    parts.append("```\n")
    parts.append(f"课件文件名：{filename}")
    parts.append(f"课件内容：\n{extracted_text[:12000]}")

    return "\n".join(parts)

def parse_ai_response(response: str) -> dict:
    """解析 AI 返回的内容"""
    result = {"formulas": [], "points": []}
    try:
        # 尝试提取 FORMULA_CARD 部分
        formula_match = re.search(
            r"===FORMULA_CARD===\s*(.*?)(?=\n===KEY_POINTS===|$)",
            response, re.DOTALL
        )
        if formula_match:
            formula_text = formula_match.group(1).strip()
            # 尝试找 JSON 数组
            json_match = re.search(r"\[[\s\S]*\]", formula_text)
            if json_match:
                try:
                    result["formulas"] = json.loads(json_match.group(0))
                except json.JSONDecodeError:
                    pass

        # 提取 KEY_POINTS 部分
        key_match = re.search(
            r"===KEY_POINTS===\s*(.*?)$",
            response, re.DOTALL
        )
        if key_match:
            key_text = key_match.group(1).strip()
            json_match = re.search(r"\[[\s\S]*\]", key_text)
            if json_match:
                try:
                    result["points"] = json.loads(json_match.group(0))
                except json.JSONDecodeError:
                    pass
    except Exception:
        pass
    return result

# ── 内容相关性检测 ──────────────────────────────────────

def check_relevance(merged_text: str) -> dict:
    """用 AI 快速检测内容是否与机器学习/计算机专业课复习相关。
    返回 {"is_relevant": bool, "reason": str}"""
    import requests as req

    sample = merged_text[:3000]  # 前 3000 字符足够判断

    prompt = f"""请判断以下文本内容是否属于「机器学习 / 计算机专业期末复习资料」范畴。

判断标准：
- 包含机器学习概念、算法、数学公式、模型训练 → 相关
- 包含计算机专业课知识点（深度学习、神经网络、NLP、数据挖掘、CV 等）→ 相关
- 不相关：招标投标、营销广告、新闻、非技术类内容

请严格用 JSON 回复，只输出 JSON：
{{"is_relevant": true或false, "reason": "一句话说明"}}

---文本---
{sample}"""

    payload = {
        "model": DEFAULT_MODEL,
        "messages": [
            {"role": "system", "content": "你是内容分类器。只输出 JSON。"},
            {"role": "user", "content": prompt}
        ],
        "temperature": 0.1,
        "max_tokens": 256
    }

    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {DEFAULT_API_KEY}"
    }

    try:
        resp = req.post(DEFAULT_API_URL, json=payload, headers=headers, timeout=(15, 30))
        if resp.status_code == 200:
            data = resp.json()
            content = data["choices"][0]["message"]["content"]
            json_match = re.search(r'\{.*\}', content, re.DOTALL)
            if json_match:
                result = json.loads(json_match.group(0))
                return {
                    "is_relevant": result.get("is_relevant", True),
                    "reason": result.get("reason", "")
                }
        return {"is_relevant": True, "reason": "检测异常，已放行"}
    except Exception:
        return {"is_relevant": True, "reason": "网络异常，已放行"}

# ── 硬编码默认 API 配置 ──────────────────────────────
DEFAULT_API_KEY = os.environ.get("DEEPSEEK_API_KEY", "")
DEFAULT_API_URL = "https://api.deepseek.com/v1/chat/completions"
DEFAULT_MODEL = "deepseek-chat"

def call_ai_api(prompt: str, api_config: dict) -> dict:
    """调用 AI API（用 requests 库，自动重试 + 长超时）"""
    import requests as req

    api_key = api_config.get("api_key", "").strip() or DEFAULT_API_KEY
    # URL 和模型由后端硬编码，不再从前端读取（防止 localStorage 缓存了旧值导致 404）
    api_url = DEFAULT_API_URL
    model = DEFAULT_MODEL

    if not api_key:
        return {"error": "请先配置 API Key", "formulas": [], "points": [],
                "raw": "⚠️ 未配置 API Key。"}

    payload = {
        "model": model,
        "messages": [
            {"role": "system", "content": "你是一位机器学习助教，擅长从课件中提取公式和重点。请严格按照要求的 JSON 格式输出。"},
            {"role": "user", "content": prompt}
        ],
        "temperature": 0.3,
        "max_tokens": 4096
    }

    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {api_key}"
    }

    # 最多重试 2 次，超时 180 秒
    for attempt in range(3):
        try:
            resp = req.post(
                api_url,
                json=payload,
                headers=headers,
                timeout=(30, 180)  # (connect_timeout, read_timeout)
            )
            if resp.status_code != 200:
                err_detail = resp.text[:500]
                return {
                    "error": f"API 返回 {resp.status_code}",
                    "formulas": [], "points": [],
                    "raw": f"❌ HTTP {resp.status_code}\n{err_detail}"
                }
            data = resp.json()
            content = data["choices"][0]["message"]["content"]
            parsed = parse_ai_response(content)
            parsed["raw"] = content
            return parsed
        except req.exceptions.ConnectTimeout:
            if attempt < 2:
                continue  # 连接超时自动重试
            return {"error": "连接超时，请检查网络或重试", "formulas": [], "points": [],
                    "raw": "❌ 无法连接到 API 服务器，请检查：\n1. 网络是否正常\n2. API 地址是否正确\n3. 是否需要科学上网"}
        except req.exceptions.ReadTimeout:
            if attempt < 2:
                continue  # 读取超时自动重试
            return {"error": "响应超时，文本可能太长，请拆分课件或重试", "formulas": [], "points": [],
                    "raw": "❌ API 响应超时（超过 3 分钟），可能课件内容太长，建议拆分上传。"}
        except Exception as e:
            if attempt < 2:
                continue
            return {"error": f"API 调用失败: {str(e)}", "formulas": [], "points": [],
                    "raw": f"❌ 调用失败（已重试 3 次）: {str(e)}"}
    return {"error": "多次重试均失败", "formulas": [], "points": [],
            "raw": "❌ 多次重试均失败，请稍后再试。"}

# ── 路由 ────────────────────────────────────────────────

@app.route("/")
def index():
    return render_template("index.html")

@app.route("/api/upload", methods=["POST"])
def upload():
    """上传文件并提取文本 — 支持多文件、ZIP 压缩包"""
    uploaded_files = request.files.getlist("file")
    if not uploaded_files or all(f.filename == "" for f in uploaded_files):
        return jsonify({"error": "没有文件"}), 400

    files_info = []
    warnings = []

    for file in uploaded_files:
        if file.filename == "":
            continue
        if not allowed_file(file.filename):
            warnings.append(f"不支持的文件格式: {file.filename}")
            continue

        # 使用 safe_basename 保留中文等 Unicode 字符
        filename = safe_basename(file.filename)
        filepath = UPLOAD_DIR / filename
        # 避免重名覆盖
        counter = 1
        stem, real_ext = os.path.splitext(filename)
        while filepath.exists():
            filepath = UPLOAD_DIR / f"{stem}_{counter}{real_ext}"
            counter += 1
        file.save(str(filepath))

        ext = get_ext(filename)

        # ZIP 解压处理
        if ext == "zip":
            zip_paths, zip_warnings = extract_zip(str(filepath))
            warnings.extend(zip_warnings)
            if not zip_paths:
                warnings.append(f"ZIP 包 \"{file.filename}\" 内未找到支持的课件文件 (PDF/PPTX/图片)")
            for zf_path in zip_paths:
                try:
                    result = extract_text(zf_path)
                    result["chars"] = len(result["text"])
                    files_info.append(result)
                except Exception as e:
                    files_info.append({"filename": os.path.basename(zf_path), "text": "", "type": "?" , "chars": 0, "error": str(e)})
        else:
            try:
                result = extract_text(str(filepath))
                result["chars"] = len(result["text"])
                files_info.append(result)
            except Exception as e:
                files_info.append({"filename": filename, "text": "", "type": "?", "chars": 0, "error": str(e)})

    if not files_info:
        return jsonify({
            "error": "未识别到有效文件。支持: PDF, PPT/PPTX, 图片, ZIP",
            "warnings": warnings
        }), 400

    # 合并文本（供单文件处理场景使用）
    merged_text = "\n\n══════ 文件分隔线 ══════\n\n".join(
        f"【文件 {i+1}: {f['filename']}】\n{f['text']}" for i, f in enumerate(files_info)
    )
    total_chars = sum(f.get("chars", 0) for f in files_info)

    # ── 内容相关性检测 ──
    relevance = {"is_relevant": True, "reason": ""}
    merged_text_for_check = "\n\n".join(f["text"][:500] for f in files_info)
    if merged_text_for_check.strip():
        relevance = check_relevance(merged_text_for_check)

    return jsonify({
        "files": files_info,
        "merged_text": merged_text,
        "total_chars": total_chars,
        "file_count": len(files_info),
        "warnings": warnings,
        "relevance": relevance
    })

@app.route("/api/process", methods=["POST"])
def process():
    """调用 AI 处理提取的文本（单文件模式，使用合并后的文本）"""
    data = request.get_json()
    if not data:
        return jsonify({"error": "缺少数据"}), 400

    text = data.get("text", "") or data.get("merged_text", "")
    file_type = data.get("type", "未知")
    filename = data.get("filename", "课件")
    api_config = data.get("api_config", {})
    mastery = data.get("mastery", "intermediate")
    urgency = data.get("urgency", None)

    if not text.strip():
        return jsonify({"error": "文本为空，请先上传文件"}), 400

    # 多文件时显示数量信息
    file_count = data.get("file_count", 1)
    if file_count > 1:
        filename = f"{filename} 等 {file_count} 份课件"

    prompt = build_prompt(text, file_type, filename, mastery, urgency)

    # 直接调用 AI
    result = call_ai_api(prompt, api_config)
    result["prompt"] = prompt
    return jsonify(result)

def build_unified_prompt(files_info: list, mastery: str = "intermediate", urgency: str = None, more_detail: bool = False) -> str:
    """构建统一提示词：同时输出公式卡、考点和思维导图，根据掌握程度和紧迫程度调整风格"""
    file_list = "\n".join(f"{i+1}. {f['filename']} ({f['type']})" for i, f in enumerate(files_info))
    is_multi = len(files_info) > 1

    # 根据掌握程度调整输出风格
    mastery_instruction = {
        "beginner": (
            "【输出风格：零基础模式】\n"
            "- 每个公式都要用大白话解释含义，不要默认用户懂任何术语\n"
            "- 考点解释要详细，给出具体例子\n"
            "- 思维导图节点文字要通俗易懂\n"
            "- 给出「这个公式是干嘛的」的一句话说明\n"
        ),
        "intermediate": (
            "【输出风格：进阶模式】\n"
            "- 公式解释简洁但完整，侧重使用场景和易错点\n"
            "- 考点突出考试重点和常见陷阱\n"
            "- 思维导图结构清晰，节点标注重点\n"
        ),
        "god": (
            "【输出风格：大佬模式】\n"
            "- 公式只给 LaTeX 和极简含义，不废话\n"
            "- 考点只列高频精华，skip 基础概念\n"
            "- 思维导图只保留核心结构和跨知识点关联\n"
            "- 可以适量使用英文术语\n"
        ),
    }.get(mastery, "【输出风格：进阶模式】\n- 适中详细程度\n")

    # 根据时间紧迫程度调整策略
    urgency_instruction = {
        "rush": (
            "【复习策略：急速模式】\n"
            "- 只提取最高优先级的公式（与 ★★★ 考点相关的）\n"
            "- 每个公式/考点用最精炼的语言，直击要害\n"
            "- 可以跳过详细解释，给「秒背版」\n"
            "- 思维导图只保留核心结构\n"
        ),
        "relaxed": (
            "【复习策略：从容模式】\n"
            "- 详细完整，每个公式给使用场景和推导思路\n"
            "- 考点配合例题思路或直观解释\n"
            "- 思维导图包含完整知识结构\n"
        ),
        "giveup": (
            "【复习策略：佛系陪伴模式】\n"
            "- 语气轻松幽默，适当使用网络梗和 emoji\n"
            "- 公式解释用最接地气的方式，帮你建立直觉\n"
            "- 考点标注「背这个就够了的保底内容」\n"
            "- 思维导图节点文字可以活泼一点\n"
        ),
    }.get(urgency, "")

    latex_example = r"$J(\theta)=\frac{1}{2m}\sum(h_\theta(x)-y)^2$"
    json_example_f = r'[{"formula": "LaTeX", "name": "名称", "meaning": "一句话含义", "source": "文件名"}]'
    json_example_p = r'[{"topic": "知识点", "summary": "核心思想", "priority": "★★★/★★/★", "pitfall": "易错提醒", "source": "文件名"}]'

    # 构建源文件映射给 AI 参考
    source_map = "\n".join(f"  - 文件名 \"{f['filename']}\" 对应第 {i+1} 份课件内容" for i, f in enumerate(files_info))

    parts = []
    parts.append("你是一位计算机专业机器学习助教。请根据以下课件内容，同时输出三样东西：\n")
    parts.append(mastery_instruction)
    if more_detail:
        parts.append(
            "\n【输出风格：超详细模式】\n"
            "- 每个公式给出完整的推导过程（从定义出发，一步一步推导）\n"
            "- 每个考点给出具体的例题或应用场景\n"
            "- 思维导图每个节点都加上详细的描述文字\n"
            "- 公式数量尽量多，不要遗漏任何重要公式\n"
        )
    if urgency_instruction:
        parts.append(urgency_instruction)
    parts.append("【一、公式速查卡】（最多 30 条，尽量多提取）")
    parts.append(f"- 提取所有数学公式，用 LaTeX 格式（如 {latex_example}）")
    parts.append("- 每个公式标注名称和含义")
    parts.append("- source 字段标记该公式出自哪个文件（见下方文件列表）\n")
    parts.append("【二、重点考点整理】")
    parts.append("- 列出知识点，标注重点程度（★★★/★★/★）和易错提醒")
    parts.append("- source 字段标记该考点出自哪个文件\n")
    mindmap_type = "跨文件知识关系" if is_multi else "章节知识结构"
    parts.append("【三、{mindmap_type}思维导图】".format(mindmap_type=mindmap_type))

    if is_multi:
        parts.append("- 分析各文件间的内在关联（前置依赖、共同主题、相互引用）")
        parts.append("- 根节点用课程主题")
        parts.append("- 二级节点按文件拆分，三级节点列出核心知识点")
    else:
        parts.append("- 分析这份课件的章节结构")
        parts.append("- 根节点用课程主题")
        parts.append("- 二级节点按章节拆分，三级节点列出核心知识点")

    parts.append("- 单独建一个 Shared_HighFreq 节点列出高频考点\n")
    parts.append("- 用 Mermaid mindmap 语法输出\n")
    parts.append("- 【绝对禁止】节点标签里不要写文件原名（如 xxx.pdf），用简短中文代替！\n")
    parts.append("- 【重要】所有节点统一用纯双引号字符串格式：\"节点名称\"\n")
    parts.append("- 【重要】禁止使用 ID[Label] 语法（如 Chap2[\"xxx\"]），直接用 \"xxx\" 即可\n")
    parts.append("- 【重要】只有根节点例外，用 root((\"课程主题\")) 格式\n")
    parts.append("- 【重要】禁止在节点里使用 emoji（如 📌📄 等），用纯文字\n")
    parts.append("- 【重要】LaTeX 公式不要写进思维导图，用纯文字描述\n")
    parts.append("- 【重要】缩进用空格：2空格=二级，4空格=三级，6空格=四级\n")

    parts.append("【输出格式严格按以下分隔】")
    parts.append("```")
    parts.append("===FORMULA_CARD===")
    parts.append(json_example_f)
    parts.append("===KEY_POINTS===")
    parts.append(json_example_p)
    parts.append("===MERMAID_MINDMAP===")
    parts.append("```mermaid")
    parts.append("mindmap")
    parts.append("  root((\"机器学习\"))")
    if is_multi:
        parts.append('    "模型评估与选择"')
        parts.append('      "性能度量"')
        parts.append('        "查准率与召回率"')
        parts.append('      "模型比较"')
        parts.append('    "决策树"')
        parts.append('      "信息增益"')
        parts.append('        "ID3算法"')
        parts.append('    "高频考点"')
        parts.append('      "过拟合"')
    else:
        parts.append('    "章节名称"')
        parts.append('      "知识点1"')
        parts.append('        "子概念"')
        parts.append('      "知识点2"')
        parts.append('    "另一章节"')
        parts.append('      "知识点3"')
        parts.append('    "高频考点"')
        parts.append('      "核心概念"')
    parts.append("```")
    parts.append("```\n")
    parts.append(f"课件文件列表：{file_list}")
    parts.append(f"\n【来源标注规则】每个公式和考点的 source 字段必须用该文件对应的文件名原文：\n{source_map}")
    parts.append("")

    # 附上每个文件的提取文本
    char_limit = 6000 if is_multi else 10000
    for i, f in enumerate(files_info):
        parts.append(f"\n===== 文件 {i+1}: {f['filename']} ({f['type']}) =====")
        parts.append(f["text"][:char_limit])

    return "\n".join(parts)

def parse_unified_response(response: str) -> dict:
    """解析统一输出：公式卡、考点、思维导图"""
    result = {"formulas": [], "points": [], "mermaid": "", "summary": ""}

    # 提取 FORMULA_CARD
    m = re.search(r"===FORMULA_CARD===\s*(.*?)(?=\n===KEY_POINTS===|$)", response, re.DOTALL)
    if m:
        jm = re.search(r"\[[\s\S]*\]", m.group(1))
        if jm:
            try:
                result["formulas"] = json.loads(jm.group(0))
            except json.JSONDecodeError:
                pass

    # 提取 KEY_POINTS
    m = re.search(r"===KEY_POINTS===\s*(.*?)(?=\n===MERMAID_MINDMAP===|$)", response, re.DOTALL)
    if m:
        jm = re.search(r"\[[\s\S]*\]", m.group(1))
        if jm:
            try:
                result["points"] = json.loads(jm.group(0))
            except json.JSONDecodeError:
                pass

    # 提取 MERMAID_MINDMAP
    mermaid_match = re.search(r"```mermaid\s*\n(.*?)```", response, re.DOTALL)
    if mermaid_match:
        raw = mermaid_match.group(1).strip()
        # 去掉外层可能有额外的 mindmap 关键字重复
        result["mermaid"] = raw

    # 简要总结（取 Mermaid 块后的文字）
    parts = re.split(r"```mermaid.*?```", response, flags=re.DOTALL)
    if len(parts) > 1:
        result["summary"] = parts[-1].strip()[:500]

    if not result["mermaid"]:
        # 回退：尝试取 ===MERMAID_MINDMAP=== 后的所有内容
        m = re.search(r"===MERMAID_MINDMAP===\s*(.*?)$", response, re.DOTALL)
        if m:
            raw = m.group(1).strip()
            raw = re.sub(r"^```mermaid\s*", "", raw)
            raw = re.sub(r"```\s*$", "", raw)
            result["mermaid"] = raw.strip()

    return result

@app.route("/api/process-unified", methods=["POST"])
def process_unified():
    """统一解析：一次返回公式卡 + 考点 + 思维导图（单文件也生成思维导图）"""
    data = request.get_json()
    if not data:
        return jsonify({"error": "缺少数据"}), 400

    files_info = data.get("files", [])
    api_config = data.get("api_config", {})
    mastery = data.get("mastery", "intermediate")
    urgency = data.get("urgency", None)
    more_detail = data.get("more_detail", False)

    if not files_info:
        return jsonify({"error": "没有已提取的文件"}), 400

    prompt = build_unified_prompt(files_info, mastery, urgency, more_detail)

    import requests as req
    api_key = api_config.get("api_key", "").strip() or DEFAULT_API_KEY

    payload = {
        "model": DEFAULT_MODEL,
        "messages": [
            {"role": "system", "content": "你是机器学习助教。请严格按照 ===SECTION=== 格式输出三部分内容，每部分只输出一次。"},
            {"role": "user", "content": prompt}
        ],
        "temperature": 0.3,
        "max_tokens": 12000
    }

    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {api_key}"
    }

    for attempt in range(3):
        try:
            resp = req.post(DEFAULT_API_URL, json=payload, headers=headers, timeout=(30, 300))
            if resp.status_code != 200:
                return jsonify({"error": f"API 返回 {resp.status_code}", "raw": resp.text[:500]})
            rd = resp.json()
            content = rd["choices"][0]["message"]["content"]
            result = parse_unified_response(content)
            result["raw"] = content
            return jsonify(result)
        except req.exceptions.Timeout:
            if attempt < 2:
                continue
            return jsonify({"error": "解析超时，请重试"})
        except Exception as e:
            if attempt < 2:
                continue
            return jsonify({"error": f"解析失败: {str(e)}"})

    return jsonify({"error": "多次重试均失败"})

@app.route("/api/generate-doc", methods=["POST"])
def generate_doc():
    """生成最终速查文档（HTML 格式）"""
    data = request.get_json()
    formulas = data.get("formulas", [])
    points = data.get("points", [])
    filename = data.get("filename", "课件")
    extracted_text = data.get("extracted_text", "")

    # 清理公式中的 $ 符号，避免模板 $$...$$ 包裹后嵌套冲突
    for f in formulas:
        if "formula" in f and f["formula"]:
            f["formula"] = f["formula"].replace("$", "").strip()

    html = render_template(
        "cheatsheet.html",
        formulas=formulas,
        points=points,
        filename=filename,
        now=datetime.now().strftime("%Y-%m-%d %H:%M"),
        text_preview=extracted_text[:2000] if extracted_text and not formulas else ""
    )

    # 保存文件
    doc_name = f"速查卡_{filename}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.html"
    doc_path = OUTPUT_DIR / doc_name
    doc_path.write_text(html, encoding="utf-8")

    return jsonify({
        "url": f"/output/{doc_name}",
        "filename": doc_name,
        "path": str(doc_path)
    })

@app.route("/output/<path:filename>")
def serve_output(filename):
    return send_from_directory(str(OUTPUT_DIR), filename)


# ── 跨文件分析 + 思维导图 ────────────────────────────────

def build_cross_analysis_prompt(files_info: list) -> str:
    """构建跨文件关联分析提示词，要求 AI 输出 Mermaid mindmap"""
    file_list = "\n".join(f"{i+1}. {f['filename']} ({f['type']})" for i, f in enumerate(files_info))

    parts = []
    parts.append("你是一位计算机专业机器学习课程助教。请你分析以下多份课件文件之间的内在关联。\n")
    parts.append("【已上传的文件清单】")
    parts.append(file_list + "\n")
    parts.append("【你需要做的】")
    parts.append("1. 识别这些文件覆盖的课程主题（根节点）")
    parts.append("2. 按文件/章节拆分为二级节点")
    parts.append("3. 在每个文件下列出其核心知识点（三级节点）")
    parts.append("4. 标注知识点之间的前置依赖和关联关系")
    parts.append("5. 找出跨文件的共同主题和高频考点")
    parts.append("6. 用 Mermaid mindmap 格式输出完整的思维导图\n")
    parts.append("【输出格式 — 严格按以下 Mermaid 语法】")
    parts.append("```mermaid")
    parts.append("mindmap")
    parts.append("  root((课程主题))")
    parts.append("    chap1[章节名称]")
    parts.append("      知识点 A")
    parts.append("        ::icon(fa fa-star)")
    parts.append("        子概念 a1")
    parts.append("        子概念 a2")
    parts.append("      知识点 B")
    parts.append("        前置依赖: 概念X")
    parts.append("    chap2[章节名称]")
    parts.append("      知识点 C")
    parts.append("        关联: 知识点 A")
    parts.append("    Shared[📌 跨章节高频考点]")
    parts.append("      核心概念 1")
    parts.append("      核心概念 2")
    parts.append("```\n")
    parts.append("【规则】")
    parts.append("- 你必须输出有效的 Mermaid mindmap 代码")
    parts.append('- 根节点使用课程主题名称（如「机器学习」）')
    parts.append("- 使用 [方括号] 的节点语法标注章节名")
    parts.append('- 用「前置依赖」和「关联」描述知识点间的关系')
    parts.append("- 单独建一个 Shared 节点展示跨文件高频考点")
    parts.append("- 仅输出 Mermaid 代码块，不要加额外解释\n")

    # 附上每个文件的提取文本
    for i, f in enumerate(files_info):
        parts.append(f"\n===== 文件 {i+1}: {f['filename']} ({f['type']}) =====")
        parts.append(f["text"][:5000])  # 每个文件限 5000 字符避免超长

    return "\n".join(parts)


def parse_mermaid_response(response: str) -> dict:
    """从 AI 响应中提取 Mermaid 代码和结构化数据"""
    result = {"mermaid": "", "summary": "", "topics": []}

    # 提取 ```mermaid ... ``` 代码块
    mermaid_match = re.search(r'```mermaid\s*\n(.*?)```', response, re.DOTALL)
    if mermaid_match:
        result["mermaid"] = mermaid_match.group(1).strip()

    # 提取简要文字总结（Mermaid 块前后的纯文本）
    summary_parts = []
    for part in re.split(r'```mermaid.*?```', response, flags=re.DOTALL):
        clean = part.strip()
        if clean:
            summary_parts.append(clean)
    result["summary"] = "\n\n".join(summary_parts[:3])  # 限 3 段

    return result


@app.route("/api/cross-analyze", methods=["POST"])
def cross_analyze():
    """跨文件关联分析，输出思维导图（Mermaid mindmap）"""
    data = request.get_json()
    if not data:
        return jsonify({"error": "缺少数据"}), 400

    files_info = data.get("files", [])
    api_config = data.get("api_config", {})

    if len(files_info) < 2:
        return jsonify({"error": "至少需要 2 个文件才能进行跨文件关联分析"}), 400

    prompt = build_cross_analysis_prompt(files_info)

    # 用相同的 AI 接口，但 max_tokens 要大一些（思维导图可能很长）
    import requests as req
    api_key = api_config.get("api_key", "").strip() or DEFAULT_API_KEY

    payload = {
        "model": DEFAULT_MODEL,
        "messages": [
            {"role": "system", "content": "你是一位机器学习助教，擅长分析课程知识体系。请输出有效的 Mermaid mindmap 代码，不要加额外解释。"},
            {"role": "user", "content": prompt}
        ],
        "temperature": 0.3,
        "max_tokens": 8192  # mindmap 可能较长，加大 token 限制
    }

    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {api_key}"
    }

    for attempt in range(3):
        try:
            resp = req.post(
                DEFAULT_API_URL,
                json=payload,
                headers=headers,
                timeout=(30, 300)  # mindmap 生成可能较慢
            )
            if resp.status_code != 200:
                return jsonify({
                    "error": f"API 返回 {resp.status_code}",
                    "mermaid": "", "summary": "",
                    "raw": f"❌ HTTP {resp.status_code}\n{resp.text[:500]}"
                })

            data = resp.json()
            content = data["choices"][0]["message"]["content"]
            parsed = parse_mermaid_response(content)
            parsed["raw"] = content

            # 同时提取公式和考点（复用已有解析逻辑）
            extra = parse_ai_response(content)
            parsed["formulas"] = extra.get("formulas", [])
            parsed["points"] = extra.get("points", [])

            return jsonify(parsed)

        except req.exceptions.Timeout:
            if attempt < 2:
                continue
            return jsonify({"error": "跨文件分析超时，请重试", "mermaid": "", "summary": ""})
        except Exception as e:
            if attempt < 2:
                continue
            return jsonify({"error": f"分析失败: {str(e)}", "mermaid": "", "summary": ""})

    return jsonify({"error": "多次重试均失败", "mermaid": "", "summary": ""})


# ── 统计 API（雷达图数据）───────────────────────────────
@app.route("/api/stats/get", methods=["GET"])
def get_stats():
    """返回当前所有选项的选择次数"""
    return jsonify(_load_stats())

@app.route("/api/stats/increment", methods=["POST"])
def stats_increment():
    """递增对应选项的计数"""
    data = request.get_json(force=True, silent=True) or {}
    mastery = data.get("mastery", "")
    urgency = data.get("urgency", "")
    result = increment_stats(mastery, urgency)
    return jsonify(result)

# ── 额外路由（挂科文档 + 本地ZIP分析）───
try:
    from extra_routes import register_routes
    register_routes(app)
    print('✅ extra_routes 已加载')
except Exception as _extra_err:
    print(f'警告: 加载 extra_routes 失败: {_extra_err}')

if __name__ == "__main__":
    print("\n" + "=" * 56)
    print("  📚 ML 速查工坊  v1.0")
    print("  课件 → AI 解析 → 公式卡 + 重点整理")
    print("=" * 56)
    print(f"  访问: http://localhost:8080")
    print(f"  上传目录: {UPLOAD_DIR}")
    print(f"  输出目录: {OUTPUT_DIR}")
    print("=" * 56 + "\n")
    app.run(debug=True, host="0.0.0.0", port=8080)
