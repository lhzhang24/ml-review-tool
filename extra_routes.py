"""
extra_routes.py - 挂科文档 + 本地ZIP分析路由
由 app.py 通过 register_routes(app) 注册
"""
import os
import re
import json
import tempfile
import urllib.request
from flask import jsonify

DEEPSEEK_API_KEY = os.environ.get("DEEPSEEK_API_KEY", "")
DEEPSEEK_API_URL = "https://api.deepseek.com/v1/chat/completions"
DEEPSEEK_MODEL = "deepseek-chat"


def call_ai_api(prompt, max_tokens=4096):
    payload = {
        "model": DEEPSEEK_MODEL,
        "messages": [
            {"role": "system", "content": "你是机器学习助教，擅长总结课件内容。"},
            {"role": "user", "content": prompt}
        ],
        "temperature": 0.3,
        "max_tokens": max_tokens
    }
    data = json.dumps(payload).encode("utf-8")
    req = urllib.request.Request(
        DEEPSEEK_API_URL,
        data=data,
        headers={
            "Content-Type": "application/json",
            "Authorization": "Bearer " + DEEPSEEK_API_KEY
        },
        method="POST"
    )
    try:
        with urllib.request.urlopen(req, timeout=180) as resp:
            result = json.loads(resp.read().decode("utf-8"))
            return result["choices"][0]["message"]["content"]
    except Exception as e:
        return "[API调用失败: " + str(e) + "]"


def register_routes(app):

    @app.route("/api/hang-ke-doc", methods=["GET"])
    def hang_ke_doc():
        html = """<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="UTF-8">
<title>考试挂科了应该怎么办？</title>
<style>
  body { font-family: -apple-system, 'PingFang SC', sans-serif; max-width: 800px; margin: 40px auto; padding: 20px; line-height: 1.8; color: #1a1a2e; }
  h1 { color: #e63946; border-bottom: 3px solid #e63946; padding-bottom: 10px; }
  .tip { background: #f1faee; border-left: 4px solid #457b9d; padding: 12px 16px; margin: 10px 0; border-radius: 0 8px 8px 0; }
  .step { background: #fff; border: 1px solid #ddd; border-radius: 8px; padding: 16px; margin: 12px 0; }
  .step h3 { margin-top: 0; color: #1d3557; }
  .btn-back { display: inline-block; margin-top: 20px; padding: 10px 24px; background: #457b9d; color: #fff; border: none; border-radius: 8px; font-size: 15px; cursor: pointer; text-decoration: none; }
  .btn-back:hover { background: #1d3557; }
  .btn-save { background: #2a9d8f; }
  .btn-save:hover { background: #2e7d32; }
  #saveme { margin-top: 20px; padding: 20px; background: #e8f5e9; border-radius: 8px; border: 1px solid #a5d6a7; display: none; }
</style>
</head>
<body>
<h1>&#x1F4CB; 考试挂科了应该怎么办？</h1>

<div class="tip" style="font-size:17px; font-weight:700; color:#e63946;">
  &#x26A0;&#xFE0F; 别急，挂科不会被枪毙！
</div>
<div class="tip">
  <span style="font-size:1.2em;">&#x1F4A1;</span> 挂科 &ne; 世界末日。大学挂科很常见，关键是知道接下来怎么处理。
</div>

<div class="step">
<h3>Step 1：确认挂科事实</h3>
<p>先去教务系统查分，确认是否真的挂了。有时候是系统延迟或者老师还没录成绩。</p>
<p>如果确实挂了 &rarr; 看 Step 2。</p>
</div>

<div class="step">
<h3>Step 2：了解补考/重修政策</h3>
<p>每所学校政策不同，通常有两种情况：</p>
<ul>
  <li><b>补考</b>：下学期开学前安排补考，通过后可以拿到学分（但绩点可能按及格算）</li>
  <li><b>重修</b>：跟下一届一起再上一次课，重新考试</li>
</ul>
<p>&#x1F449; 去学院教务办问清楚：补考时间、重修报名时间、是否影响保研/毕业。</p>
</div>

<div class="step">
<h3>Step 3：制定复习计划（如果需要补考/重修）</h3>
<p>挂科说明这块知识有漏洞，补考/重修是最后一次补救机会：</p>
<ul>
  <li>&#x1F4DA; 把课件重新过一遍（用本工具的「公式速查卡」功能！）</li>
  <li>&#x1F4DD; 把往年题刷一遍，重点看错题</li>
  <li>&#x1F64B; 主动找老师答疑，展示你的态度（老师可能会酌情给分）</li>
</ul>
</div>

<div class="step">
<h3>Step 4（最重要）：调整心态</h3>
<p>挂科后最容易陷入「破罐破摔」的恶性循环：</p>
<ul>
  <li>挂了一门 &rarr; 心情差 &rarr; 其他课也不好好学 &rarr; 挂更多</li>
</ul>
<p><b>正确做法</b>：把这次挂科当作一次「early warning」，及时调整学习方法和时间分配。</p>
</div>

<div class="tip" style="background:#fdf2e9; border-color:#f4a261;">
  <span style="font-size:1.2em;">&#x1F434;</span> <b>佛系挂科语录</b>：挂科是对这门课的不够热爱，不是对你能力的否定。调整方向，下次再来！
</div>

<br>
<a class="btn-back" href="/" onclick="localStorage.setItem('mltool_skip_onboard','1');">&#x2190; 返回复习工具</a>
<br><br>
<button class="btn-back btn-save" onclick="document.getElementById('saveme').style.display='block'; this.style.display='none';">
  &#x1F69C; 我觉得我还能再救一救
</button>

<div id="saveme">
  <h3 style="margin-top:0; color:#2e7d32;">&#x1F69C; 抢救方案</h3>
  <p>好！既然还想救，就认真做以下几件事：</p>
  <ul>
    <li>① 用本工具上传你的课件，生成「公式速查卡」—— 这是最高效的复习方式</li>
    <li>② 把挂的那门课的课件全部重新过一遍，重点看 &#x2B50;&#x2B50;&#x2B50; 考点</li>
    <li>③ 找同学要一下往年的考试题，刷题是最快的提分方式</li>
    <li>④ 如果补考/重修还是没把握，考虑找助教或者学霸辅导</li>
  </ul>
  <p><b>你现在就去上传课件，开始复习。别再拖了。</b></p>
  <a class="btn-back" href="/" onclick="localStorage.setItem('mltool_skip_onboard','1');" style="background:#2e7d32;">好，去上传课件 &rarr;</a>
</div>

</body>
</html>"""
        return html


    @app.route("/api/analyze-local-zip", methods=["GET"])
    def analyze_local_zip():
        zip_path = "/Users/main/Downloads/课件.zip"
        if not os.path.exists(zip_path):
            return jsonify({"error": "文件不存在: " + zip_path})

        tmpdir = tempfile.mkdtemp(prefix="ml_zip_")
        try:
            with zipfile.ZipFile(zip_path, "r") as zf:
                for name in zf.namelist():
                    if name.endswith("/"):
                        continue
                    basename = os.path.basename(name)
                    if "__MACOSX" in name or basename.startswith("._"):
                        continue
                    if basename.startswith("."):
                        continue
                    safe = re.sub(r'[\x00-\x1f\x7f\\/*?"<>|]', '_', basename)
                    dest = os.path.join(tmpdir, safe)
                    with zf.open(name) as src, open(dest, "wb") as dst:
                        dst.write(src.read())
        except Exception as e:
            return jsonify({"error": "解压失败: " + str(e)})

        # 导入提取函数（从主模块）
        import sys
        main_module = sys.modules[__name__.split('.')[0] if '.' in __name__ else __name__]
        
        all_texts = {}
        for fname in os.listdir(tmpdir):
            fpath = os.path.join(tmpdir, fname)
            if not os.path.isfile(fpath):
                continue
            ext = ""
            if "." in fname:
                ext = fname.rsplit(".", 1)[1].lower()
            # 简单提取：读取前几页
            text = ""
            try:
                if ext == "pdf":
                    try:
                        from pdfplumber import open as pdf_open
                        with pdf_open(fpath) as pdf:
                            for i, page in enumerate(pdf.pages[:5]):
                                t = page.extract_text()
                                if t:
                                    text += t + "\n"
                    except Exception:
                        pass
                elif ext in ("pptx", "ppt"):
                    try:
                        from pptx import Presentation
                        prs = Presentation(fpath)
                        for i, slide in enumerate(prs.slides[:8]):
                            for shape in slide.shapes:
                                if shape.has_text_frame:
                                    for para in shape.text_frame.paragraphs:
                                        if para.text.strip():
                                            text += para.text + " "
                    except Exception:
                        pass
            except Exception:
                pass
            if text and len(text) > 100:
                all_texts[fname] = text[:8000]

        if not all_texts:
            return jsonify({"error": "没有提取到任何文本内容，可能是扫描件"})

        # 按章节分组
        chapters = {}
        for fname, text in all_texts.items():
            m = re.search(r"chap(\d+)", fname, re.IGNORECASE)
            if m:
                ch = "第" + m.group(1) + "章"
            elif "复习" in fname:
                ch = "期末复习"
            elif "模糊" in fname or "FCM" in fname or "fuzzy" in fname.lower():
                ch = "模糊C均值"
            else:
                ch = "其他"
            if ch not in chapters:
                chapters[ch] = []
            chapters[ch].append({"filename": fname, "text": text})

        # 逐章总结
        results = {}
        for ch, files in chapters.items():
            combined = "\n".join("=== " + f["filename"] + " ===\n" + f["text"][:6000] for f in files)
            prompt = ("以下是机器学习课程「" + ch + "」的课件内容，请做一份完整的复习总结，包括：\n" +
                      "1. 本章核心知识点（★★★/★★/★）\n" +
                      "2. 重要公式（LaTeX格式）及含义\n" +
                      "3. 常见考点和易错点\n" +
                      "4. 复习建议\n\n" +
                      "课件内容：\n" + combined[:20000] + "\n\n请用Markdown格式输出。")
            summary = call_ai_api(prompt, max_tokens=4096)
            results[ch] = summary

        # 总复习大纲
        total_prompt = "以下是机器学习各章节的总结，请生成一份期末总复习大纲：\n\n"
        for ch, summary in sorted(results.items()):
            total_prompt += "### " + ch + " ###\n" + summary[:2000] + "\n\n"
        total_prompt += "\n请输出：1. 知识框架 2. 公式总表 3. 期末重点（按考试概率排序） 4. 复习时间分配建议"

        total_summary = call_ai_api(total_prompt, max_tokens=8192)

        return jsonify({
            "chapters": results,
            "total_summary": total_summary,
            "files_analyzed": list(all_texts.keys())
        })



    @app.route("/api/zip-summary-doc", methods=["GET"])
    def zip_summary_doc():
        """返回课件.zip的AI总结文档HTML，格式正确，内容精简"""
        import markdown
        import re

        summary_path = "/tmp/zip_chapter_summaries.json"
        if not os.path.exists(summary_path):
            return """<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="UTF-8">
<title>正在生成总结...</title>
<style>
  body { font-family: -apple-system, 'PingFang SC', sans-serif;
         display: flex; justify-content: center; align-items: center;
         height: 100vh; margin: 0; background: #f5f5f5; }
  .box { text-align: center; padding: 40px; background: white;
          border-radius: 12px; box-shadow: 0 2px 12px rgba(0,0,0,0.1); }
  .spinner { font-size: 40px; animation: spin 1s linear infinite; }
  @keyframes spin { from { transform: rotate(0deg); } to { transform: rotate(360deg); } }
</style>
</head>
<body>
<div class="box">
  <div class="spinner">⏳</div>
  <h2>正在AI生成课件总结...</h2>
  <p>预计需要5-10分钟，请稍后刷新页面</p>
  <a href="/api/zip-summary-doc">刷新</a> | <a href="/">返回</a>
</div>
</body>
</html>"""

        with open(summary_path, "r", encoding="utf-8") as f:
            data = json.load(f)

        chapters = {k: v for k, v in data.items() if k != "总复习大纲"}
        total_summary = data.get("总复习大纲", "")

        def convert_md(text):
            """将Markdown转为HTML，并处理LaTeX公式"""
            if not text:
                return ""
            import re
            # 先保护公式不被 markdown 处理
            placeholders = {}
            cnt = [0]

            def protect_inline(m):
                key = f"@@LATEX_INLINE_{cnt[0]}@@"
                placeholders[key] = m.group(0)
                cnt[0] += 1
                return key

            def protect_block(m):
                key = f"@@LATEX_BLOCK_{cnt[0]}@@"
                placeholders[key] = m.group(0)
                cnt[0] += 1
                return key

            # 保护公式块 $$...$$
            text_protected = re.sub(r'\$\$(.*?)\$\$', protect_block, text, flags=re.DOTALL)
            # 保护行内公式 $...$
            text_protected = re.sub(r'(?<!\$)\$(?!\$)(.+?)(?<!\$)\$(?!\$)', protect_inline, text_protected)

            # 转换 Markdown 为 HTML（使用 markdown 函数）
            html = markdown.markdown(text_protected, extensions=["extra", "toc"])

            # 恢复公式
            for key, val in placeholders.items():
                if key.startswith("@@LATEX_BLOCK_"):
                    formula = val[2:-2].strip()
                    # 用 KaTeX 渲染：直接输出 $$ formula $$，让前端 auto-render 处理
                    html = html.replace(key, f'<div class="formula-block">$$ {formula} $$</div>')
                else:
                    formula = val[1:-1].strip()
                    html = html.replace(key, f'<span class="latex-inline">$ {formula} $</span>')

            return html

        # === 构建完整 HTML ===
        html_parts = []
        html_parts.append("""<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>机器学习课件总结 - 急速复习版</title>

<!-- KaTeX -->
<link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/katex@0.16.9/dist/katex.min.css">
<script src="https://cdn.jsdelivr.net/npm/katex@0.16.9/dist/katex.min.js"></script>
<script src="https://cdn.jsdelivr.net/npm/katex@0.16.9/dist/contrib/auto-render.min.js"></script>

<style>
  * { box-sizing: border-box; }
  body { font-family: -apple-system, 'PingFang SC', 'Microsoft YaHei', sans-serif;
         max-width: 900px; margin: 0 auto; padding: 20px;
         line-height: 1.9; color: #1e293b; background: #f8fafc; }

  .header { text-align: center; margin-bottom: 24px; padding: 20px;
            background: linear-gradient(135deg, #6366f1 0%, #8b5cf6 100%);
            border-radius: 14px; color: white; }
  .header h1 { margin: 0 0 6px; font-size: 1.5em; color: white; border: none; }
  .header p { margin: 0; opacity: 0.9; font-size: 0.95em; }

  .toc { background: white; border: 1px solid #e0e7ff; border-radius: 10px;
         padding: 14px 18px; margin: 16px 0; box-shadow: 0 1px 3px rgba(0,0,0,0.04); }
  .toc-title { font-weight: 700; font-size: 1em; color: #4338ca; margin-bottom: 8px; }
  .toc a { display: block; padding: 4px 0; color: #4f46e5; text-decoration: none;
           font-size: 0.92em; }
  .toc a:hover { color: #6366f1; text-decoration: underline; }

  .chapter { background: white; border: 1px solid #e2e8f0; border-radius: 10px;
             padding: 18px 22px; margin: 14px 0;
             box-shadow: 0 1px 3px rgba(0,0,0,0.04); }
  .chapter h2 { color: #4f46e5; margin: 0 0 12px; font-size: 1.2em;
                border-left: 4px solid #6366f1; padding-left: 10px; }
  .chapter h3 { color: #334155; margin: 12px 0 6px; font-size: 1.05em; }

  .formula-block { background: #f8fafc; border: 1px solid #cbd5e1; border-radius: 8px;
                   padding: 12px 16px; margin: 10px 0; text-align: center;
                   overflow-x: auto; }

  ul { padding-left: 20px; }
  li { margin: 5px 0; line-height: 1.8; }

  .btn-back { display: inline-block; margin-top: 20px; padding: 10px 24px;
              background: linear-gradient(135deg, #6366f1, #8b5cf6);
              color: #fff; border: none; border-radius: 8px;
              font-size: 0.95em; cursor: pointer; text-decoration: none; }
  .btn-back:hover { opacity: 0.9; }

  .katex-display { margin: 12px 0 !important; }
  .katex { font-size: 1.08em !important; }
</style>

<script>
window.addEventListener("DOMContentLoaded", function() {
    renderMathInElement(document.body, {
        delimiters: [
            {left: "$$", right: "$$", display: true},
            {left: "$", right: "$", display: false}
        ],
        throwOnError: false,
        trust: true
    });
});
</script>
</head>
<body>
""")

        # 头部
        html_parts.append(f"""<div class="header">
<h1>📚 机器学习课件总结</h1>
<p>急速复习版 · 基于 {len(chapters)} 个章节 · AI 自动生成</p>
</div>""")

        # 目录
        html_parts.append('<div class="toc"><div class="toc-title">📋 目录导航</div>')
        for i, ch in enumerate(sorted(chapters.keys()), 1):
            html_parts.append(f'<a href="#ch{i}">{i}. {ch}</a>')
        if total_summary:
            html_parts.append('<a href="#total">📋 总复习大纲</a>')
        html_parts.append('</div>\n')

        # 各章节
        for i, (ch, summary) in enumerate(sorted(chapters.items()), 1):
            html_parts.append(f'<div class="chapter" id="ch{i}">\n<h2>{ch}</h2>\n')
            html_parts.append(convert_md(summary))
            html_parts.append('</div>\n')

        # 总复习大纲
        if total_summary:
            html_parts.append('<div class="chapter" id="total">\n<h2>📋 总复习大纲</h2>\n')
            html_parts.append(convert_md(total_summary))
            html_parts.append('</div>\n')

        html_parts.append('<div style="text-align:center;"><a class="btn-back" href="/">← 返回复习工具</a></div>\n')
        html_parts.append('</body>\n</html>')

        return '\n'.join(html_parts)

