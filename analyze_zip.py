"""
analyze_zip.py - 分析课件.zip：提取所有课件文本，调用 AI 做全面总结
用法: python3 analyze_zip.py
"""
import os
import re
import json
from pathlib import Path

ZIP_DIR = Path(__file__).resolve().parent / "uploads" / "zip_analysis"
OUTPUT_DIR = Path(__file__).resolve().parent / "outputs"
OUTPUT_DIR.mkdir(exist_ok=True)

# ── 文本提取 ──────────────────────────────────────────
try:
    from PyPDF2 import PdfReader
except ImportError:
    PdfReader = None

try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


def get_ext(filename):
    return filename.rsplit(".", 1)[1].lower() if "." in filename else ""


def extract_pdf_text(filepath, max_pages=5):
    text = ""
    if pdfplumber:
        try:
            with pdfplumber.open(filepath) as pdf:
                for i, page in enumerate(pdf.pages):
                    if i >= max_pages:
                        break
                    t = page.extract_text()
                    if t:
                        text += f"[第{i+1}页]\n{t}\n"
            return text.strip()
        except Exception:
            pass
    if PdfReader:
        try:
            reader = PdfReader(filepath)
            for i in range(min(max_pages, len(reader.pages))):
                t = reader.pages[i].extract_text()
                if t:
                    text += f"[第{i+1}页]\n{t}\n"
        except Exception:
            pass
    return text.strip()


def extract_pptx_text(filepath, max_slides=8):
    if not Presentation:
        return "[python-pptx 未安装]"
    try:
        prs = Presentation(filepath)
        text = ""
        for i, slide in enumerate(prs.slides):
            if i >= max_slides:
                break
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            if texts:
                text += f"[第{i+1}页]\n" + "\n".join(texts) + "\n\n"
        return text.strip()
    except Exception as e:
        return f"[PPTX解析失败: {e}]"


def extract_all():
    results = []
    files = sorted(ZIP_DIR.iterdir(), key=lambda f: f.name)
    for fpath in files:
        ext = get_ext(fpath.name)
        print(f"  提取: {fpath.name} ({ext.upper()})")
        if ext == "pdf":
            text = extract_pdf_text(str(fpath), max_pages=5)
        elif ext in ("pptx", "ppt"):
            text = extract_pptx_text(str(fpath), max_slides=8)
        else:
            text = ""
        results.append({
            "filename": fpath.name,
            "ext": ext,
            "text": text[:6000],
            "chars": len(text)
        })
    return results


# ── AI 总结 ───────────────────────────────────────────
DEEPSEEK_API_KEY = os.environ.get("DEEPSEEK_API_KEY", "")
DEEPSEEK_API_URL = "https://api.deepseek.com/v1/chat/completions"
DEEPSEEK_MODEL = "deepseek-chat"

import urllib.request


def call_deepseek(prompt, max_tokens=4096):
    payload = {
        "model": DEEPSEEK_MODEL,
        "messages": [
            {"role": "system", "content": "你是一位机器学习助教，擅长总结课件内容，输出清晰的结构化笔记。"},
            {"role": "user", "content": prompt}
        ],
        "temperature": 0.3,
        "max_tokens": max_tokens
    }
    data = json.dumps(payload).encode("utf-8")
    req = urllib.request.Request(
        DEEPSEEK_API_URL,
        data=data,
        headers={
            "Content-Type": "application/json",
            "Authorization": f"Bearer {DEEPSEEK_API_KEY}"
        },
        method="POST"
    )
    try:
        with urllib.request.urlopen(req, timeout=180) as resp:
            result = json.loads(resp.read().decode("utf-8"))
            return result["choices"][0]["message"]["content"]
    except Exception as e:
        return f"[API调用失败: {e}]"


def summarize_all(files_info):
    chapters = {}
    for f in files_info:
        name = f["filename"]
        m = re.search(r'chap(\d+)', name, re.IGNORECASE)
        if m:
            ch = m.group(1)
        elif '复习' in name:
            ch = "review"
        elif '模糊' in name:
            ch = "fuzzy_cmeans"
        else:
            ch = "other"
        if ch not in chapters:
            chapters[ch] = []
        chapters[ch].append(f)

    print(f"\n共 {len(files_info)} 个文件，分 {len(chapters)} 组")
    summaries = {}
    for ch, flist in chapters.items():
        print(f"  总结章节 {ch} ({len(flist)} 个文件)...")
        combined = ""
        for f in flist:
            combined += f"\n===== {f['filename']} =====\n{f['text']}\n"
        if len(combined) > 24000:
            combined = combined[:24000] + "\n[内容过长已截断]"

        prompt = f"""以下是机器学习课程第{ch}章的课件内容（可能包含多个文件/多个部分）。
请做一份完整的章节总结，包括：
1. 本章核心知识点列表（带★★★/★★/★优先级）
2. 重要公式（LaTeX格式）及含义
3. 常见考点和易错点
4. 各小节之间的关系

课件内容：
{combined}

请输出结构清晰的 Markdown 格式总结。"""

        summary = call_deepseek(prompt, max_tokens=4096)
        summaries[ch] = summary
        print(f"  完成，总结长度: {len(summary)} 字符")

    # 最终汇总
    print("\n  生成最终汇总...")
    final_prompt = """以下是机器学习课程各章节的AI总结，请生成一份完整的课程总复习大纲，包括：
1. 课程整体知识框架
2. 各章节核心内容概览
3. 公式总表（按章节排列）
4. 期末复习重点（按考试概率排序）
5. 复习建议

各章节总结如下：
"""
    for ch, summary in sorted(summaries.items()):
        final_prompt += f"\n### 第{ch}章 ###\n{summary}\n"

    final_summary = call_deepseek(final_prompt, max_tokens=8192)
    return summaries, final_summary


def main():
    print("=" * 60)
    print("  课件.zip 内容分析")
    print("=" * 60)

    print("\n[1/3] 扫描文件...")
    files_info = extract_all()
    print(f"  完成，共 {len(files_info)} 个文件")
    total = sum(f["chars"] for f in files_info)
    print(f"  提取文本总计: {total:,} 字符")

    print("\n[2/3] 调用 AI 总结各章节...")
    summaries, final_summary = summarize_all(files_info)

    print("\n[3/3] 保存结果...")
    for ch, summary in summaries.items():
        out_path = OUTPUT_DIR / f"章节总结_第{ch}章.md"
        out_path.write_text(f"# 第{ch}章 总结\n\n{summary}", encoding="utf-8")
        print(f"  已保存: {out_path.name}")

    final_path = OUTPUT_DIR / "课程总复习大纲.md"
    final_path.write_text(f"# 机器学习课程总复习大纲\n\n{final_summary}", encoding="utf-8")
    print(f"  已保存: {final_path.name}")

    all_text_path = OUTPUT_DIR / "全部课件原始文本.txt"
    with open(all_text_path, "w", encoding="utf-8") as f:
        for fi in files_info:
            f.write(f"\n{'='*60}\n")
            f.write(f"文件: {fi['filename']}\n")
            f.write(f"字符数: {fi['chars']}\n")
            f.write(f"{'='*60}\n")
            f.write(fi["text"])
            f.write("\n\n")
    print(f"  已保存: {all_text_path.name}")

    print(f"\n{'='*60}")
    print(f"  完成！所有结果已保存到:")
    print(f"  {OUTPUT_DIR}")
    print(f"{'='*60}")


if __name__ == "__main__":
    main()
